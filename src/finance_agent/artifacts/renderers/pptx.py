"""pptx 渲染器：ArtifactSpec（slide blocks）→ 决策框架演示文稿。

页数与内容由 spec 决定；渲染器负责一套确定性的版式系统：
- 16:9 画幅，全部页面在空白版式上自绘（不依赖模板占位符——默认模板的
  占位符排版正是旧版"字体乱、溢出、挤成一团"的根源）；
- 浅色主题（与 Web UI / HTML 产物同一色板）、统一字号层级与页脚；
- 防溢出是渲染器的责任而非 LLM 自律：要点页/表格页超容量自动分"（续）"页，
  长标题自动降字号。
"""

from __future__ import annotations

import io
from collections.abc import Iterator, Mapping
from datetime import UTC, datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from finance_agent.artifacts.spec import ArtifactSpec, SlideBlock
from finance_agent.provenance import Evidence
from finance_agent.skills.loader import SkillInfo

SUPPORTED_BLOCKS = {"slide"}

# ---- 画幅与网格（16:9） ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.66)
CONTENT_W = Emu(SLIDE_W - 2 * MARGIN)
BODY_TOP = Inches(1.5)
BODY_BOTTOM = Inches(6.9)

# ---- 色板（与 webapp 浅色主题一致） ----
_TEXT = RGBColor(0x1F, 0x27, 0x33)
_MUTED = RGBColor(0x5B, 0x6B, 0x81)
_ACCENT = RGBColor(0x2A, 0x63, 0xE8)
_LINE = RGBColor(0xD7, 0xDE, 0xE9)
_PANEL = RGBColor(0xF2, 0xF5, 0xFA)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_FONT_LATIN = "Helvetica Neue"
_FONT_EA = "PingFang SC"

# ---- 容量（超出即自动分页，防溢出） ----
MAX_BULLETS_PER_SLIDE = 6
MAX_TABLE_ROWS_PER_SLIDE = 8
MAX_EVIDENCE_PER_SLIDE = 12


class RenderError(RuntimeError):
    pass


# ---------- 基础绘制 ----------

def _style(run, size: float, *, bold: bool = False, color: RGBColor = _TEXT) -> None:
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    font.name = _FONT_LATIN
    # 中文字体需写 a:ea（python-pptx 的 font.name 只设 latin）
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", _FONT_EA)


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def _para(frame, text: str, size: float, *, bold=False, color=_TEXT,
          align=PP_ALIGN.LEFT, space_after: float = 0, first=False):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    if space_after:
        para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    _style(run, size, bold=bold, color=color)
    return para


def _bar(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _set_bullet_char(para) -> None:
    """真正的项目符号（悬挂缩进 + buChar），换行后正文对齐而非顶到符号下。"""
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "228600")   # 0.25 英寸
    pPr.set("indent", "-228600")
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "•"}))


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _chunks(seq: list, size: int) -> Iterator[list]:
    for i in range(0, len(seq), size):
        yield seq[i:i + size]


def _cont_title(title: str, index: int) -> str:
    return title if index == 0 else f"{title}（续）"


# ---------- 页面骨架 ----------

def _content_header(slide, title: str) -> None:
    size = 22 if len(title) <= 26 else 18   # 长标题降字号而非溢出
    frame = _textbox(slide, MARGIN, Inches(0.46), CONTENT_W, Inches(0.78))
    frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    _para(frame, title, size, bold=True, first=True)
    _bar(slide, MARGIN, Inches(1.3), Inches(1.05), Inches(0.045), _ACCENT)


def _footer(slide, page_no: int, note: str = "") -> None:
    if note:
        frame = _textbox(slide, MARGIN, Inches(7.06), Emu(CONTENT_W - Inches(0.6)), Inches(0.3))
        _para(frame, note, 8.5, color=_MUTED, first=True)
    right = _textbox(slide, Emu(SLIDE_W - MARGIN - Inches(0.6)), Inches(7.06),
                     Inches(0.6), Inches(0.3))
    _para(right, str(page_no), 9, color=_MUTED, align=PP_ALIGN.RIGHT, first=True)


def _block_note(block: SlideBlock) -> str:
    parts = []
    if block.evidence_refs:
        parts.append(f"溯源：{'、'.join(block.evidence_refs)}")
    if block.notes:
        parts.append(block.notes)
    return "  ｜  ".join(parts)


# ---------- 版式 ----------

def _add_title_slide(prs, block: SlideBlock, _page) -> None:
    slide = _blank(prs)
    _bar(slide, MARGIN, Inches(2.32), Inches(1.05), Inches(0.06), _ACCENT)
    title_frame = _textbox(slide, MARGIN, Inches(2.62), CONTENT_W, Inches(1.9))
    _para(title_frame, block.title, 34 if len(block.title) <= 30 else 28,
          bold=True, first=True)
    if block.subtitle:
        sub = _textbox(slide, MARGIN, Inches(4.55), CONTENT_W, Inches(1.0))
        _para(sub, block.subtitle, 16, color=_MUTED, first=True)
    meta = _textbox(slide, MARGIN, Inches(6.7), CONTENT_W, Inches(0.35))
    _para(meta, f"finance-agent 研究产物 · {datetime.now(UTC):%Y-%m-%d}",
          10, color=_MUTED, first=True)


def _add_section_slide(prs, block: SlideBlock, page) -> None:
    slide = _blank(prs)
    _bar(slide, Emu((SLIDE_W - Inches(0.9)) // 2), Inches(2.95),
         Inches(0.9), Inches(0.055), _ACCENT)
    frame = _textbox(slide, MARGIN, Inches(3.25), CONTENT_W, Inches(1.1))
    _para(frame, block.title, 30, bold=True, align=PP_ALIGN.CENTER, first=True)
    if block.subtitle:
        sub = _textbox(slide, MARGIN, Inches(4.4), CONTENT_W, Inches(0.8))
        _para(sub, block.subtitle, 15, color=_MUTED, align=PP_ALIGN.CENTER, first=True)
    _footer(slide, page(), _block_note(block))


def _add_bullets_slide(prs, block: SlideBlock, page) -> None:
    bullets = block.bullets or [""]
    for i, chunk in enumerate(_chunks(bullets, MAX_BULLETS_PER_SLIDE)):
        slide = _blank(prs)
        _content_header(slide, _cont_title(block.title, i))
        size = 18 if len(chunk) <= 4 else 16
        frame = _textbox(slide, MARGIN, BODY_TOP, CONTENT_W,
                         Emu(BODY_BOTTOM - BODY_TOP))
        for j, bullet in enumerate(chunk):
            para = _para(frame, bullet, size, space_after=10, first=(j == 0))
            para.line_spacing = 1.2
            _set_bullet_char(para)
        _footer(slide, page(), _block_note(block))


def _add_table_slide(prs, block: SlideBlock, page) -> None:
    if not block.table_headers:
        raise RenderError(f"table 版式的 slide（{block.title}）缺少 table_headers")
    cols = len(block.table_headers)
    body_size = 12 if cols <= 4 else 11 if cols <= 5 else 10
    for i, chunk in enumerate(_chunks(block.table_rows or [], MAX_TABLE_ROWS_PER_SLIDE)):
        slide = _blank(prs)
        _content_header(slide, _cont_title(block.title, i))
        n_rows = len(chunk) + 1
        shape = slide.shapes.add_table(
            n_rows, cols, MARGIN, BODY_TOP, CONTENT_W,
            Emu(Inches(0.5) + Inches(0.42) * len(chunk)),
        )
        table = shape.table
        table.first_row = False       # 关掉模板自带的深蓝配色/条纹
        table.horz_banding = False
        if cols >= 2:                 # 首列（维度/标签列）略宽
            first_w = int(CONTENT_W * 1.3 / (cols + 0.3))
            other_w = int((CONTENT_W - first_w) / (cols - 1))
            table.columns[0].width = Emu(first_w)
            for c in range(1, cols):
                table.columns[c].width = Emu(other_w)
        table.rows[0].height = Inches(0.5)
        for c, header in enumerate(block.table_headers):
            _fill_cell(table.cell(0, c), header, body_size, bold=True,
                       color=_WHITE, fill=_ACCENT)
        for r, row in enumerate(chunk, start=1):
            table.rows[r].height = Inches(0.42)
            fill = _PANEL if r % 2 == 0 else _WHITE
            for c in range(cols):
                text = row[c] if c < len(row) else ""
                _fill_cell(table.cell(r, c), text, body_size,
                           bold=(c == 0), fill=fill)
        _footer(slide, page(), _block_note(block))


def _fill_cell(cell, text: str, size: float, *, bold=False,
               color: RGBColor = _TEXT, fill: RGBColor | None = None) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill is not None else _WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.1)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    frame = cell.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    _style(run, size, bold=bold, color=color)


def _evidence_slides(prs, evidence: Mapping[str, Evidence], page) -> None:
    items = list(evidence.values())
    lines = [
        f"{ev.id}［{ev.kind}］{ev.source_url}（{ev.fetched_at}）" for ev in items
    ] or ["（本演示未登记 evidence）"]
    for i, chunk in enumerate(_chunks(lines, MAX_EVIDENCE_PER_SLIDE)):
        slide = _blank(prs)
        _content_header(slide, _cont_title("数据来源与溯源清单", i))
        frame = _textbox(slide, MARGIN, BODY_TOP, CONTENT_W,
                         Emu(BODY_BOTTOM - BODY_TOP))
        for j, line in enumerate(chunk):
            para = _para(frame, line, 10.5, color=_MUTED, space_after=6, first=(j == 0))
            para.line_spacing = 1.1
        _footer(slide, page())


_LAYOUT_DISPATCH = {
    "title": _add_title_slide,
    "section": _add_section_slide,
    "bullets": _add_bullets_slide,
    "table": _add_table_slide,
}


def render_pptx(
    spec: ArtifactSpec,
    *,
    datasets: Mapping = {},
    evidence: Mapping[str, Evidence],
    skill: SkillInfo,
) -> bytes:
    if spec.kind != "pptx":
        raise RenderError(f"kind 不匹配：期望 pptx，得到 {spec.kind}")
    unsupported = {b.type for b in spec.blocks} - SUPPORTED_BLOCKS
    if unsupported:
        raise RenderError(f"pptx 渲染器不支持 block 类型：{sorted(unsupported)}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page_counter = {"n": 1}

    def page() -> int:
        page_counter["n"] += 1
        return page_counter["n"]

    slides = [b for b in spec.blocks if isinstance(b, SlideBlock)]
    if slides[0].layout != "title":
        # 自动补齐封面，保证任何 spec 都有标题页
        _add_title_slide(prs, SlideBlock(
            layout="title", title=spec.title,
            subtitle=spec.subtitle or datetime.now(UTC).strftime("%Y-%m-%d"),
        ), page)
    for block in slides:
        _LAYOUT_DISPATCH[block.layout](prs, block, page)
    _evidence_slides(prs, evidence, page)

    prs.core_properties.title = spec.title
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
