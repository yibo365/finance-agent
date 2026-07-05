"""office 三渲染器单测：渲染 → 用对应库重新打开 → 断言结构与公式。"""

from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pandas as pd
import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from finance_agent.artifacts.spec import ArtifactSpec
from finance_agent.workspace import Workspace

_W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


@pytest.fixture()
def ws(tmp_path):
    workspace = Workspace.create(tmp_path / "outputs", "s-20260703-office")
    dates_5d = pd.date_range("2024-01-01", periods=60, freq="B").strftime("%Y-%m-%d")
    dates_7d = pd.date_range("2024-01-01", periods=80, freq="D").strftime("%Y-%m-%d")

    def make(dates, base):
        closes = [base + i * 0.5 for i in range(len(dates))]
        return pd.DataFrame({
            "date": dates, "open": closes, "high": [c * 1.01 for c in closes],
            "low": [c * 0.99 for c in closes], "close": closes,
            "volume": [1000] * len(dates),
        })

    workspace.store_dataset("ds-gold", make(dates_5d, 2000), ticker="GC=F")
    workspace.store_dataset("ds-btc", make(dates_7d, 40000), ticker="BTC-USD")
    workspace.evidence.record("market_data", source_url="https://example.com/gold", excerpt="60 行")
    workspace.save_evidence()
    return workspace


def _docx_xml(path, member: str):
    with ZipFile(path) as zf:
        return ET.fromstring(zf.read(member))


def _zip_member_text(path, member: str) -> str:
    with ZipFile(path) as zf:
        return zf.read(member).decode("utf-8")


def _style_rfonts(styles, style_id: str):
    style = styles.find(f".//{_W_NS}style[@{_W_NS}styleId='{style_id}']")
    assert style is not None, style_id
    return style.find(f"./{_W_NS}rPr/{_W_NS}rFonts")


def test_xlsx_backtest_formulas_and_sheets(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "gold-btc-backtest",
        "kind": "xlsx",
        "title": "黄金 vs 比特币回测底稿",
        "blocks": [
            {"type": "heading", "text": "口径声明"},
            {"type": "narrative", "text": "多资产按共同交易日对齐。",
             "evidence_refs": ["ev-s-20260703-office-1"]},
            {"type": "data_sheet", "sheet_name": "黄金", "data_ref": "ds-gold", "ticker": "GC=F"},
            {"type": "data_sheet", "sheet_name": "比特币", "data_ref": "ds-btc", "ticker": "BTC-USD"},
            {"type": "metrics_sheet", "data_refs": ["ds-gold", "ds-btc"],
             "labels": ["黄金", "比特币"], "rolling_window": 20},
        ],
    })
    version = ws.render_artifact(spec)
    wb = load_workbook(ws.dir / version.file)
    assert {"说明", "黄金", "比特币", "参数", "对齐", "指标", "汇总", "年度收益",
            "图表", "溯源"} <= set(wb.sheetnames)
    # 参数化：窗口值在参数 sheet，指标公式引用它（OFFSET 联动）
    assert wb["参数"]["B2"].value == 20
    vol_formula = wb["指标"]["E30"].value  # 黄金滚动波动列
    assert isinstance(vol_formula, str) and "OFFSET" in vol_formula and "'参数'!$B$2" in vol_formula
    # 日收益、回撤是公式非数值
    assert str(wb["指标"]["B3"].value).startswith("=")
    assert "MAX(" in wb["指标"]["D10"].value
    # 汇总：核心指标齐全且全为公式（真实事故：模型只能靠 table 写"公式"占位）
    summary = {row[0].value: row[1].value for row in wb["汇总"].iter_rows(min_row=2)}
    for name in ("区间总收益", "年化收益率（CAGR）", "年化波动率（全样本）",
                 "夏普比率（rf=0）", "最大回撤", "正收益天数占比", "日收益相关系数"):
        assert name in summary and str(summary[name]).startswith("="), name
    assert "CORREL" in summary["日收益相关系数"]
    assert "POWER" in summary["年化收益率（CAGR）"]

    # 年度收益：按对齐数据真实年份逐年 + 全周期 + 差额列，公式引用对齐 sheet
    annual = wb["年度收益"]
    aligned_last = wb["对齐"].max_row
    year_rows = {row[0].value: row for row in annual.iter_rows(min_row=2)
                 if row[0].value and str(row[0].value) != "注：年度收益以上一年最后共同交易日收盘为基准（首年度以区间首日为基准）；首末年度可能为区间内的部分年度。"}
    assert "2024" in year_rows and "全周期" in year_rows   # 夹具数据全在 2024
    y2024 = year_rows["2024"]
    assert y2024[1].value == f"='对齐'!B{aligned_last}/'对齐'!B2-1"   # 首年以区间首日为基准
    assert y2024[3].value.startswith("=C") and "-B" in y2024[3].value  # 差额列
    full_row = year_rows["全周期"]
    assert full_row[2].value == f"='对齐'!C{aligned_last}/'对齐'!C2-1"
    # 对齐 sheet 行数 = 共同交易日（工作日 60 天 ⊂ 自然日 80 天中的工作日）
    assert wb["对齐"].max_row - 1 <= 60
    # 溯源 sheet 有记录
    assert wb["溯源"].max_row >= 2


def test_xlsx_rejects_placeholder_table_cells(ws):
    # 真实事故：年度收益/核心指标两个 table 全是字面量"公式"，产物看似成功实为空壳
    spec = ArtifactSpec.model_validate({
        "artifact_id": "placeholder-backtest", "kind": "xlsx", "title": "占位符",
        "blocks": [
            {"type": "data_sheet", "sheet_name": "黄金", "data_ref": "ds-gold", "ticker": "GC=F"},
            {"type": "table", "caption": "年度收益对比", "headers": ["年份", "收益"],
             "rows": [["2024", "公式"]]},
        ],
    })
    with pytest.raises(Exception, match="占位符"):
        ws.render_artifact(spec)


def test_xlsx_requires_data_sheet(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "empty-backtest", "kind": "xlsx", "title": "空",
        "blocks": [{"type": "heading", "text": "x"}],
    })
    with pytest.raises(Exception, match="data_sheet"):
        ws.render_artifact(spec)


def test_pptx_slides_and_evidence_page(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "gold-btc-deck",
        "kind": "pptx",
        "title": "黄金与比特币：避险与抗通胀比较",
        "blocks": [
            {"type": "slide", "layout": "title", "title": "黄金与比特币：避险与抗通胀比较",
             "subtitle": "决策框架 v1"},
            {"type": "slide", "layout": "section", "title": "一、比较框架"},
            {"type": "slide", "layout": "bullets", "title": "核心结论",
             "bullets": ["黄金危机相关性更稳定", "比特币波动率显著更高"],
             "evidence_refs": ["ev-s-20260703-office-1"]},
            {"type": "slide", "layout": "table", "title": "多维对照",
             "table_headers": ["维度", "黄金", "比特币"],
             "table_rows": [["年化波动", "低", "高"], ["历史长度", "长", "短"]]},
        ],
    })
    version = ws.render_artifact(spec)
    prs = Presentation(ws.dir / version.file)
    # 16:9 画幅
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)
    assert len(prs.slides) == 5  # 4 页 spec + 1 页自动溯源清单

    def slide_text(slide):
        return " ".join(
            run.text
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs for run in para.runs
        )

    assert "数据来源与溯源清单" in slide_text(prs.slides[4])
    # 表格页确实有表格
    table_slide = prs.slides[3]
    assert any(shape.has_table for shape in table_slide.shapes)
    # 所有形状都落在画幅内（旧版默认模板正是排版溢出的重灾区）
    for slide in prs.slides:
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= prs.slide_width
            assert shape.top + shape.height <= prs.slide_height + 1


def test_pptx_evidence_urls_are_clickable_hyperlinks(ws):
    """真实事故：PPTX 证据页只写普通 URL 文本，不能一跳打开来源。"""
    spec = ArtifactSpec.model_validate({
        "artifact_id": "hyperlink-deck",
        "kind": "pptx",
        "title": "PPTX 超链接测试",
        "blocks": [
            {"type": "slide", "layout": "title", "title": "PPTX 超链接测试"},
        ],
    })
    version = ws.render_artifact(spec)
    with ZipFile(ws.dir / version.file) as zf:
        rel_members = [
            name for name in zf.namelist()
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels")
        ]
        rels = "\n".join(zf.read(name).decode("utf-8") for name in rel_members)
    assert "relationships/hyperlink" in rels
    assert 'Target="https://example.com/gold"' in rels
    assert 'TargetMode="External"' in rels


def test_pptx_overflow_splits_into_continuation_slides(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "overflow-deck", "kind": "pptx", "title": "容量测试",
        "blocks": [
            {"type": "slide", "layout": "bullets", "title": "十条要点",
             "bullets": [f"要点 {i}" for i in range(10)]},
            {"type": "slide", "layout": "table", "title": "长表",
             "table_headers": ["维度", "值"],
             "table_rows": [[f"行{i}", str(i)] for i in range(20)]},
        ],
    })
    version = ws.render_artifact(spec)
    prs = Presentation(ws.dir / version.file)
    # 自动封面 1 + bullets 2（6+4）+ table 3（8+8+4）+ 溯源 1 = 7
    assert len(prs.slides) == 7
    texts = [
        " ".join(r.text for s in slide.shapes if s.has_text_frame
                 for p in s.text_frame.paragraphs for r in p.runs)
        for slide in prs.slides
    ]
    assert sum("（续）" in t for t in texts) == 3


def test_pptx_rejects_non_slide_blocks(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "bad-deck", "kind": "pptx", "title": "x",
        "blocks": [{"type": "heading", "text": "不该出现在PPT里"}],
    })
    with pytest.raises(Exception, match="不支持"):
        ws.render_artifact(spec)


def test_docx_structure_and_appendix(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "gold-btc-report",
        "kind": "docx",
        "title": "黄金与比特币策略研究报告",
        "subtitle": "避险/抗通胀属性比较",
        "blocks": [
            {"type": "heading", "text": "一、摘要"},
            {"type": "narrative", "text": "黄金的危机期相关性更稳定。\n\n比特币波动率更高。",
             "evidence_refs": ["ev-s-20260703-office-1"]},
            {"type": "heading", "text": "二、数据与方法"},
            {"type": "table", "caption": "数据说明", "headers": ["标的", "区间"],
             "rows": [["GC=F", "2024"], ["BTC-USD", "2024"]]},
            {"type": "changepoint_table", "changepoints": [
                {"date": "2024-02-01", "kind": "drawdown", "rule": "回撤 -16%",
                 "severity": 3, "window": ["2024-01-15", "2024-02-01"]},
            ]},
        ],
    })
    version = ws.render_artifact(spec)
    doc = Document(str(ws.dir / version.file))
    texts = [p.text for p in doc.paragraphs]
    headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert "一、摘要" in headings and "附录：溯源明细" in headings
    assert any("ev-s-20260703-office-1" in t for t in texts)   # 正文溯源标记 + 附录
    assert any("不构成投资建议" in t for t in texts)
    assert len(doc.tables) == 2  # 数据说明 + 变化点明细


def test_docx_evidence_urls_are_clickable_hyperlinks(ws):
    """真实事故：Word 附录只写普通 URL 文本，不能一跳打开来源。"""
    spec = ArtifactSpec.model_validate({
        "artifact_id": "hyperlink-report",
        "kind": "docx",
        "title": "Word 超链接测试",
        "blocks": [{"type": "narrative", "text": "正文。"}],
    })
    version = ws.render_artifact(spec)
    rels = _zip_member_text(ws.dir / version.file, "word/_rels/document.xml.rels")
    document = _zip_member_text(ws.dir / version.file, "word/document.xml")
    assert "relationships/hyperlink" in rels
    assert 'Target="https://example.com/gold"' in rels
    assert 'TargetMode="External"' in rels
    assert "<w:hyperlink" in document


def test_docx_uses_chinese_report_fonts(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "font-report",
        "kind": "docx",
        "title": "中文字体测试",
        "blocks": [
            {"type": "heading", "text": "一、摘要"},
            {"type": "narrative", "text": "正文里的中文、Ticker GC=F 和数字 2026 应该使用统一报告字体。"},
        ],
    })
    version = ws.render_artifact(spec)
    docx_path = ws.dir / version.file

    styles = _docx_xml(docx_path, "word/styles.xml")
    settings = _docx_xml(docx_path, "word/settings.xml")
    for style_id in ("Normal", "Title", "Heading1"):
        rfonts = _style_rfonts(styles, style_id)
        assert rfonts is not None, style_id
        assert rfonts.attrib[f"{_W_NS}eastAsia"] == "PingFang SC"
        assert rfonts.attrib[f"{_W_NS}ascii"] == "Helvetica Neue"
        assert rfonts.attrib[f"{_W_NS}hAnsi"] == "Helvetica Neue"

    theme_lang = settings.find(f".//{_W_NS}themeFontLang")
    assert theme_lang is not None
    assert theme_lang.attrib[f"{_W_NS}eastAsia"] == "zh-CN"


def test_docx_rejects_slide_block(ws):
    spec = ArtifactSpec.model_validate({
        "artifact_id": "bad-doc", "kind": "docx", "title": "x",
        "blocks": [{"type": "slide", "layout": "bullets", "title": "PPT block"}],
    })
    with pytest.raises(Exception, match="不支持"):
        ws.render_artifact(spec)
